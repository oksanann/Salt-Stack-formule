#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Generate PowerPoint: Salt formula generation project."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "salt-formula-generation.pptx"

ACCENT = RGBColor(0x1A, 0x56, 0xDB)
DARK = RGBColor(0x1E, 0x1E, 0x1E)
MUTED = RGBColor(0x5A, 0x5A, 0x5A)


def set_run(run, *, size=18, bold=False, color=DARK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    for p in slide.shapes.title.text_frame.paragraphs:
        for r in p.runs:
            set_run(r, size=36, bold=True, color=ACCENT)
    for p in slide.placeholders[1].text_frame.paragraphs:
        for r in p.runs:
            set_run(r, size=20, color=MUTED)


def add_section_slide(prs, section_num, title):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = f"Слайд {section_num} · {title}"
    for p in slide.shapes.title.text_frame.paragraphs:
        for r in p.runs:
            set_run(r, size=32, bold=True, color=ACCENT)


def add_bullet_slide(prs, title, bullets, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
        p.space_after = Pt(8)
        for r in p.runs:
            set_run(r, size=18)
    if note:
        p = body.add_paragraph()
        p.text = note
        p.level = 0
        p.space_before = Pt(16)
        for r in p.runs:
            set_run(r, size=14, color=MUTED)


def add_table_slide(prs, title, headers, rows, *, top=1.45, height=5.2, font_size=11):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    for p in slide.shapes.title.text_frame.paragraphs:
        for r in p.runs:
            set_run(r, size=28, bold=True, color=ACCENT)

    cols = len(headers)
    table_shape = slide.shapes.add_table(
        len(rows) + 1, cols, Inches(0.35), Inches(top), Inches(9.3), Inches(height)
    )
    table = table_shape.table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                set_run(r, size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    set_run(r, size=font_size)
    return table


HYPOTHESES = [
    (
        "H1",
        "ИИ сразу пишет все файлы «рецепта»",
        "Просим «сделай PuTTY» — получаем папку в ответе",
        "Не подошло",
        "Каждый раз другой ответ, легко забыть файл, долго проверять",
    ),
    (
        "H2",
        "ИИ пишет bash-скрипт, который собирает папку",
        "Скрипт на bash создаёт готовую формулу",
        "Не подошло",
        "Пустые файлы, ошибки в кавычках; новый скрипт на каждый продукт",
    ),
    (
        "H3",
        "ИИ пишет Python-скрипт (build_*_formula.py)",
        "python3 build_….py → готовая папка",
        "Лучше bash, но не выбрали",
        "Python стабильнее, но ИИ каждый раз пишет новый сборщик",
    ),
    (
        "H4",
        "ИИ пишет только JSON, сборку делает наш скрипт",
        "JSON → render_formula.sh → dist/",
        "Выбрали",
        "Один JSON + один проверенный скрипт → предсказуемый результат",
    ),
    (
        "H5",
        "База знаний (RAG) с правилами и примерами",
        "ИИ видит правила перед генерацией JSON",
        "Подтвердили",
        "Меньше ошибок в ОС, названиях и ссылках",
    ),
]

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Как мы научили ИИ собирать «рецепты» установки ПО",
        "Осмакс · T1 Innotech · Август 2026\n"
        "Простыми словами — без погружения в Salt",
    )

    add_bullet_slide(
        prs,
        "1. Задача",
        [
            "Осмакс удалённо устанавливает программы на Linux (Astra, ALT) и Windows.",
            "Для каждой программы нужен «рецепт» — папка с инструкциями: что установить, откуда, как удалить.",
            "Боль: ручная сборка долгая, легко ошибиться в структуре или забыть Windows.",
            "Цель: инженер описывает задачу словами → получает готовую проверенную папку.",
        ],
        note="Аналогия: формула = рецепт в поваренной книге. ИИ пишет черновик, наш скрипт оформляет в нужном формате.",
    )

    add_table_slide(
        prs,
        "2. Гипотезы и результаты проверки",
        ["№", "Идея", "Как проверяли", "Итог", "Почему"],
        HYPOTHESES,
        font_size=10,
    )

    add_bullet_slide(
        prs,
        "3. Что проверяем после сборки",
        [
            "JSON читается без ошибок (jq)",
            "Есть имя продукта, список ОС, настройки установки",
            "После сборки нет пустых файлов",
            "Структура папки совпадает с тем, что ждёт Осмакс",
        ],
        note="Вывод: bash ломался чаще всего. Python — надёжнее, но JSON + один скрипт — самый стабильный путь.",
    )

    add_bullet_slide(
        prs,
        "4. Решение, которое работает",
        [
            "ИИ заполняет одну «анкету» (JSON). Сборку делает один скрипт render_formula.sh.",
            "",
            "Цепочка:",
            "  Запрос → ИИ + база знаний (rag/) → *.json → скрипт → dist/",
            "",
            "Папки в проекте:",
            "  rag/ — инструкции для ИИ",
            "  json-formula/ — описания формул (PuTTY, браузер, безопасность…)",
            "  dist/ — готовые папки для установки",
            "",
            "Примеры: PuTTY, Яндекс Браузер, security baseline, Horizon Client.",
        ],
        note="Главное: предсказуемость. Обновили скрипт — все формулы собираются по новым правилам.",
    )

    add_bullet_slide(
        prs,
        "5. Куда двигаемся дальше",
        [
            "Песочница — быстрая проверка:",
            "  • Docker/WSL: закинул JSON → получил папку и отчёт «всё ок»",
            "  • Автопроверка всех формул в CI",
            "",
            "Приложение для пользователя:",
            "  • Поле «что установить» + «на каких ОС»",
            "  • ИИ → JSON → скрипт → zip с готовой формулой",
            "  • Отчёт: что проверили и прошло ли",
        ],
        note="Следующий шаг: обёртка вокруг уже работающего скрипта, без переписывания с нуля.",
    )

    # Closing
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Спасибо"
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(2))
    tf = tx.text_frame
    tf.text = (
        "Репозиторий: osmax-salt-formulas-rag\n"
        "Рабочий пайплайн: JSON-first + render_formula.sh"
    )
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            set_run(r, size=22, color=MUTED)

    prs.save(OUT)
    print(f"OK: {OUT}")


if __name__ == "__main__":
    main()
